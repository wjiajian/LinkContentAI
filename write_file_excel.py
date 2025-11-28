import os
import sys
import warnings
import logging

# 抑制所有PDF相关警告
logging.getLogger("pdfminer").setLevel(logging.ERROR)

import openpyxl
import docx
import tempfile
import shutil
import uuid
import re
import json
import base64
from typing import Dict, List, Tuple, Optional
from pathlib import Path
from openpyxl.utils import get_column_letter
from openai import OpenAI
# from xbot import print

# 抑制所有库的警告
for mod in ['pdfplumber', 'pdf2image', 'PIL']:
    try:
        warnings.filterwarnings('ignore', module=mod)
    except:
        pass

# 配置多模态LLM
# TODO: 请配置您的qwen-vl API信息
QWEN_VL_CONFIG = {
    "api_key": os.getenv("QWEN_V"),  # API密钥
    "base_url": "https://dashscope.aliyuncs.com/compatible-mode/v1",  # 通义千问API endpoint
    "model": "qwen-vl-plus",  # 或 qwen-vl-max
}

# 临时文件管理类
class TempFileManager:
    """管理临时文件和目录的生命周期"""
    def __init__(self):
        self.temp_dir = None
        self.used_paths = set()

    def __enter__(self):
        self.temp_dir = tempfile.mkdtemp(prefix="excel_img_proc_")
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.temp_dir and os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir, ignore_errors=True)

    def get_temp_path(self, suffix="") -> str:
        """生成唯一的临时文件路径"""
        unique_id = str(uuid.uuid4())
        filename = f"{unique_id}{suffix}"
        self.used_paths.add(filename)
        return os.path.join(self.temp_dir, filename)

# --- 模块化的内容读取区域 ---
# TODO: 这里可以添加更多的文件类型支持
# 未来若要添加对新文件类型（例如 .csv）的支持:
# 1. 编写一个新的函数 `read_csv_content(file_path)`。
# 2. 在 FILE_READERS 字典中增加一行映射：`'.csv': read_csv_content`。

def read_txt_content(file_path: str) -> str:
    """从 .txt 文件中读取内容。"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"读取 TXT 文件 '{file_path}' 时出错: {e}"

def read_docx_content(file_path: str) -> str:
    """从 .docx 文件中读取内容。"""
    try:
        doc = docx.Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return '\n'.join(full_text)
    except Exception as e:
        return f"读取 DOCX 文件 '{file_path}' 时出错: {e}"

def read_xlsx_content(file_path: str) -> str:
    """
    从 .xlsx 文件中的所有工作表读取可见的文本内容。
    """
    try:
        # 以只读模式加载工作簿，这样性能更好，且不会意外修改文件
        workbook = openpyxl.load_workbook(file_path, read_only=True)

        all_sheets_text = []

        # 遍历工作簿中的每一个工作表
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []

            # 添加工作表标题，以便区分不同工作表的内容
            sheet_text.append(f"--- 工作表: {sheet.title} ---")

            # 遍历工作表中的每一行
            for row in sheet.iter_rows():
                # 获取行中每个单元格的值，并转换为字符串，忽略空单元格
                # str(cell.value) 可以安全地处理数字、日期等不同类型
                row_values = [str(cell.value) for cell in row if cell.value is not None]

                # 如果行中有内容，则将它们用制表符连接起来
                if row_values:
                    sheet_text.append("\t".join(row_values))

            # 将当前工作表的所有文本行用换行符连接起来
            all_sheets_text.append("\n".join(sheet_text))

        # 将所有工作表的内容用两个换行符隔开，使其更清晰
        return "\n\n".join(all_sheets_text)

    except FileNotFoundError:
        return f"错误：Excel 文件未找到 '{file_path}'"
    except Exception as e:
        return f"读取 XLSX 文件 '{file_path}' 时出错: {e}"


def read_pptx_content(file_path: str) -> str:
    """
    从 .pptx 文件中读取文本内容。
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        all_slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- 幻灯片 {slide_num} ---"]

            # 提取幻灯片中的所有形状的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(f"{shape.text.strip()}\n")

            all_slides_text.append("\n".join(slide_text))

        return "\n\n".join(all_slides_text)

    except ImportError:
        return "错误：需要安装 python-pptx 库来读取PPTX文件: pip install python-pptx"
    except Exception as e:
        error_msg = str(e)
        if "Package not found" in error_msg or "未找到" in error_msg:
            return f"错误：PPTX 文件未找到 '{file_path}'"
        return f"读取 PPTX 文件 '{file_path}' 时出错: {error_msg}"


def read_xmind_content(file_path: str) -> str:
    """
    从 .xmind 文件中读取文本内容。
    使用xmindparser库将XMind文件转换为Python字典，然后提取文本内容。
    """
    try:
        from xmindparser import xmind_to_dict

        # 使用xmindparser解析XMind文件
        xmind_data = xmind_to_dict(file_path)

        # 递归提取所有主题的文本内容
        all_text = []

        def extract_text_recursive(topic_data, level=0):
            """递归提取主题文本"""
            if isinstance(topic_data, dict):
                # 提取当前主题的标题
                if 'title' in topic_data:
                    indent = "  " * level
                    title = topic_data['title']
                    if title and title.strip():
                        all_text.append(f"{indent}- {title.strip()}")

                # 处理注释
                if 'note' in topic_data and topic_data['note']:
                    indent = "  " * level
                    note = topic_data['note']
                    if note.strip():
                        all_text.append(f"{indent}  注释: {note.strip()}")

                # 处理标签
                if 'labels' in topic_data and topic_data['labels']:
                    indent = "  " * level
                    labels = topic_data['labels']
                    if labels:
                        all_text.append(f"{indent}  标签: {', '.join(labels)}")

                # 处理链接
                if 'link' in topic_data and topic_data['link']:
                    indent = "  " * level
                    link = topic_data['link']
                    if link.strip():
                        all_text.append(f"{indent}  链接: {link.strip()}")

                # 递归处理子主题
                if 'topics' in topic_data and topic_data['topics']:
                    for sub_topic in topic_data['topics']:
                        extract_text_recursive(sub_topic, level + 1)

        # 遍历所有工作表
        if isinstance(xmind_data, list):
            for sheet in xmind_data:
                if 'topic' in sheet:
                    # 添加工作表标题
                    if 'title' in sheet:
                        all_text.append(f"\n=== {sheet['title']} ===\n")
                    else:
                        all_text.append(f"\n=== 工作表 ===\n")

                    extract_text_recursive(sheet['topic'])
        elif isinstance(xmind_data, dict) and 'topic' in xmind_data:
            extract_text_recursive(xmind_data['topic'])

        return "\n".join(all_text) if all_text else "无法解析XMind文件内容"

    except ImportError:
        return "错误：需要安装 xmindparser 库来读取XMind文件: pip install xmindparser"
    except FileNotFoundError:
        return f"错误：XMind 文件未找到 '{file_path}'"
    except Exception as e:
        return f"读取 XMind 文件 '{file_path}' 时出错: {e}"


def read_pdf_content(file_path: str) -> str:
    """
    从 .pdf 文件中读取文本内容。
    """
    try:
        import pdfplumber

        # 抑制PDF字体警告
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            all_text = []
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        all_text.append(f"--- 第 {page_num} 页 ---\n{page_text}")

        return "\n\n".join(all_text)

    except ImportError:
        return "错误：需要安装 pdfplumber 库来读取PDF文件: pip install pdfplumber"
    except FileNotFoundError:
        return f"错误：PDF 文件未找到 '{file_path}'"
    except Exception as e:
        return f"读取 PDF 文件 '{file_path}' 时出错: {e}"


# --- 图片提取功能 ---
def extract_images_from_docx(docx_path: str, temp_manager: TempFileManager) -> List[str]:
    """
    从 DOCX 文件中提取所有嵌入的图片。
    返回提取的图片路径列表。
    """
    try:
        import zipfile
        import os

        image_paths = []
        docx_dir = tempfile.mkdtemp(prefix="docx_extract_")

        # DOCX 实际上是一个ZIP文件
        with zipfile.ZipFile(docx_path, 'r') as zip_ref:
            zip_ref.extractall(docx_dir)
            media_dir = os.path.join(docx_dir, "word", "media")

            if os.path.exists(media_dir):
                for filename in os.listdir(media_dir):
                    if any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']):
                        src_path = os.path.join(media_dir, filename)
                        # 复制到临时目录
                        temp_path = temp_manager.get_temp_path(suffix=f"_{filename}")
                        shutil.copy2(src_path, temp_path)
                        image_paths.append(temp_path)

        # 清理临时目录
        shutil.rmtree(docx_dir, ignore_errors=True)
        return image_paths

    except Exception as e:
        print(f"从DOCX提取图片时出错: {e}")
        return []


def extract_images_from_pdf(pdf_path: str, temp_manager: TempFileManager) -> List[str]:
    """
    从 PDF 文件中提取图片。
    返回提取的图片路径列表。
    """
    try:
        # 尝试使用 pdf2image 将PDF转换为图片
        from pdf2image import convert_from_path

        images = convert_from_path(pdf_path)
        image_paths = []

        for idx, img in enumerate(images):
            temp_path = temp_manager.get_temp_path(suffix=f"_page_{idx+1}.png")
            img.save(temp_path, 'PNG')
            image_paths.append(temp_path)

        return image_paths

    except ImportError:
        print("警告：需要安装 pdf2image 来处理PDF图片: pip install pdf2image")
        print("     还需要安装 Poppler: https://pdf2image.readthedocs.io/en/latest/installation.html")
        return []
    except Exception as e:
        print(f"从PDF提取图片时出错: {e}")
        return []


def extract_images_from_pptx(pptx_path: str, temp_manager: TempFileManager) -> List[str]:
    """
    从 PPTX 文件中提取所有嵌入的图片。
    返回提取的图片路径列表。
    """
    try:
        import zipfile
        import os

        image_paths = []
        pptx_dir = tempfile.mkdtemp(prefix="pptx_extract_")

        # PPTX 实际上是一个ZIP文件
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            zip_ref.extractall(pptx_dir)
            media_dir = os.path.join(pptx_dir, "ppt", "media")

            if os.path.exists(media_dir):
                for filename in os.listdir(media_dir):
                    if any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']):
                        src_path = os.path.join(media_dir, filename)
                        # 复制到临时目录
                        temp_path = temp_manager.get_temp_path(suffix=f"_{filename}")
                        shutil.copy2(src_path, temp_path)
                        image_paths.append(temp_path)

        # 清理临时目录
        shutil.rmtree(pptx_dir, ignore_errors=True)
        return image_paths

    except Exception as e:
        print(f"从PPTX提取图片时出错: {e}")
        return []


def extract_images_from_xmind(xmind_path: str, temp_manager: TempFileManager) -> List[str]:
    """
    从 XMind 文件中提取所有嵌入的图片。
    返回提取的图片路径列表。
    """
    try:
        import zipfile
        import os

        image_paths = []
        xmind_dir = tempfile.mkdtemp(prefix="xmind_extract_")

        try:
            # XMind 实际上是一个ZIP文件
            with zipfile.ZipFile(xmind_path, 'r') as zip_ref:
                zip_ref.extractall(xmind_dir)

                # 查找媒体文件目录（可能在不同的位置）
                media_dirs = []
                for root, dirs, files in os.walk(xmind_dir):
                    # XMind的媒体文件可能在OOML、content或media目录中
                    if any(keyword in root.lower() for keyword in ['media', 'ooml', 'images']):
                        media_dirs.append(root)

                # 如果没有找到特定目录，查找所有目录
                if not media_dirs:
                    media_dirs = [xmind_dir]

                for media_dir in media_dirs:
                    if os.path.exists(media_dir):
                        for filename in os.listdir(media_dir):
                            if any(filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.svg']):
                                src_path = os.path.join(media_dir, filename)
                                # 复制到临时目录
                                temp_path = temp_manager.get_temp_path(suffix=f"_{filename}")
                                shutil.copy2(src_path, temp_path)
                                image_paths.append(temp_path)

        finally:
            # 清理临时目录
            shutil.rmtree(xmind_dir, ignore_errors=True)

        return image_paths

    except Exception as e:
        print(f"从XMind提取图片时出错: {e}")
        return []


def extract_images_from_document(file_path: str, temp_manager: TempFileManager) -> List[str]:
    """
    从任何支持的文档中提取图片。
    """
    _, extension = os.path.splitext(file_path.lower())

    if extension == '.docx':
        return extract_images_from_docx(file_path, temp_manager)
    elif extension == '.pdf':
        return extract_images_from_pdf(file_path, temp_manager)
    elif extension == '.pptx':
        return extract_images_from_pptx(file_path, temp_manager)
    elif extension == '.xmind':
        return extract_images_from_xmind(file_path, temp_manager)
    else:
        return []


# --- 文档转Markdown功能 ---
def convert_docx_to_markdown_with_placeholders(docx_path: str, image_paths: List[str], temp_manager: TempFileManager) -> str:
    """
    将DOCX转换为带占位符的Markdown。
    通过XML解析检测图片在文档中的精确位置并插入占位符。
    """
    try:
        import zipfile
        import xml.etree.ElementTree as ET

        # 使用python-docx读取文档
        doc = docx.Document(docx_path)

        markdown_lines = []
        image_idx = 0

        # 通过XML解析来精确检测图片位置
        docx_zip = zipfile.ZipFile(docx_path)
        document_xml = docx_zip.read('word/document.xml')
        root = ET.fromstring(document_xml)

        # 定义命名空间
        ns = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
            'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture'
        }

        # 查找所有图片及其位置
        image_positions = []
        for idx, para in enumerate(root.findall('.//w:p', ns)):
            # 检查此段落是否包含图片 - 使用多种方式检测
            # 方式1: 检查wp:docPr (drawing properties)
            has_image1 = para.find('.//wp:docPr', ns) is not None
            # 方式2: 检查a:blip (bitmap image)
            has_image2 = para.find('.//a:blip', ns) is not None
            # 方式3: 检查pic:pic (picture)
            has_image3 = para.find('.//pic:pic', ns) is not None

            has_image = has_image1 or has_image2 or has_image3

            if has_image:
                image_positions.append(idx)

        docx_zip.close()

        # 使用精确插入
        for para_idx, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    markdown_lines.append(f"{'#' * int(level)} {text}\n")
                else:
                    markdown_lines.append(text + "\n")

            # 如果当前段落有图片，插入占位符
            if para_idx in image_positions and image_idx < len(image_paths):
                markdown_lines.append(f"![placeholder]({image_paths[image_idx]})\n")
                image_idx += 1

        # 如果还有剩余图片，追加到末尾
        while image_idx < len(image_paths):
            markdown_lines.append(f"![placeholder]({image_paths[image_idx]})\n")
            image_idx += 1

        return "\n".join(markdown_lines)

    except Exception as e:
        return f"转换DOCX时出错: {e}"


def convert_pdf_to_markdown_with_placeholders(pdf_path: str, image_paths: List[str]) -> str:
    """
    将PDF转换为带占位符的Markdown。
    通过页面图片检测功能检测页面中的图片位置并插入占位符。
    """
    try:
        import pdfplumber

        # 抑制PDF字体警告
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            markdown_lines = []
            image_idx = 0

            with pdfplumber.open(pdf_path) as pdf:
                page_texts = []
                page_image_counts = []

                # 提取所有页面的文本和图片信息
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    page_texts.append(page_text if page_text else "")

                    # 检测页面中的图片数量
                    image_count = 0
                    if hasattr(page, 'images') and page.images:
                        image_count = len(page.images)

                    page_image_counts.append(image_count)

                # 生成Markdown，按检测到的图片位置插入
                for page_num, (page_text, image_count) in enumerate(zip(page_texts, page_image_counts), 1):
                    markdown_lines.append(f"--- 第 {page_num} 页 ---\n")
                    if page_text:
                        markdown_lines.append(page_text)

                    # 如果检测到页面有图片，插入相应数量的占位符
                    if image_count > 0 and image_idx < len(image_paths):
                        for _ in range(image_count):
                            if image_idx < len(image_paths):
                                markdown_lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
                                image_idx += 1

                # 如果还有剩余图片，追加到最后一页
                while image_idx < len(image_paths):
                    markdown_lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
                    image_idx += 1

        return "\n\n".join(markdown_lines)

    except Exception as e:
        return f"转换PDF时出错: {e}"


def convert_pptx_to_markdown_with_placeholders(pptx_path: str, image_paths: List[str], temp_manager: TempFileManager) -> str:
    """
    将PPTX转换为带占位符的Markdown。
    按幻灯片提取文本，根据幻灯片中的图片形状插入占位符。
    """
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        markdown_lines = []
        image_idx = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            # 添加幻灯片标题
            markdown_lines.append(f"--- 幻灯片 {slide_num} ---")

            # 提取幻灯片内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 检查是否为标题
                    if shape.is_placeholder:
                        placeholder = shape.placeholder_format
                        if placeholder.type == 1:  # 标题占位符
                            markdown_lines.append(f"\n## {shape.text.strip()}\n")
                        else:
                            markdown_lines.append(f"{shape.text.strip()}\n")
                    else:
                        markdown_lines.append(f"{shape.text.strip()}\n")

            # 检查幻灯片中是否有图片形状
            slide_has_image = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    slide_has_image = True
                    break

            # 如果幻灯片有图片，插入占位符
            if slide_has_image and image_idx < len(image_paths):
                markdown_lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
                image_idx += 1

        # 如果还有剩余图片，追加到末尾
        while image_idx < len(image_paths):
            markdown_lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
            image_idx += 1

        return "\n".join(markdown_lines)

    except Exception as e:
        return f"转换PPTX时出错: {e}"


def convert_xmind_to_markdown_with_placeholders(xmind_path: str, image_paths: List[str], temp_manager: TempFileManager) -> str:
    """
    将XMind转换为带占位符的Markdown。
    使用xmindparser库将XMind文件转换为Python字典，然后格式化为Markdown。
    """
    try:
        from xmindparser import xmind_to_dict

        # 使用xmindparser解析XMind文件
        xmind_data = xmind_to_dict(xmind_path)

        markdown_lines = []
        image_idx = 0

        def extract_topic_recursive(topic_data, level=0):
            """递归提取主题结构并插入占位符"""
            lines = []

            if isinstance(topic_data, dict):
                # 生成标题（基于层级）
                if level == 0:
                    # 根主题用一级标题
                    header_prefix = "#"
                elif level == 1:
                    # 二级主题
                    header_prefix = "##"
                else:
                    # 其他层级用项目符号
                    header_prefix = None

                # 提取主题标题
                if 'title' in topic_data:
                    title = topic_data['title']
                    if title and title.strip():
                        if header_prefix:
                            lines.append(f"{header_prefix} {title.strip()}\n")
                        else:
                            lines.append(f"{'  ' * (level - 1)}- {title.strip()}\n")

                # 处理注释
                if 'note' in topic_data and topic_data['note']:
                    note = topic_data['note']
                    if note.strip():
                        lines.append(f"{'  ' * level}> {note.strip()}\n")

                # 处理标签
                if 'labels' in topic_data and topic_data['labels']:
                    labels = topic_data['labels']
                    if labels:
                        lines.append(f"{'  ' * level}标签: {', '.join(labels)}\n")

                # 处理链接
                if 'link' in topic_data and topic_data['link']:
                    link = topic_data['link']
                    if link.strip():
                        lines.append(f"{'  ' * level}链接: {link.strip()}\n")

                # 检查是否有图片标记（通过markers或其他属性判断）
                has_image_indicator = False
                if 'makers' in topic_data and topic_data['makers']:
                    # 如果有特殊的标记，可以作为图片占位符的指示器
                    markers = topic_data['makers']
                    if any('image' in marker.lower() or 'picture' in marker.lower() for marker in markers):
                        has_image_indicator = True

                # 在主题内容后插入占位符（如果需要）
                if has_image_indicator and image_idx < len(image_paths):
                    lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
                    image_idx += 1

                # 递归处理子主题
                if 'topics' in topic_data and topic_data['topics']:
                    for sub_topic in topic_data['topics']:
                        sub_lines = extract_topic_recursive(sub_topic, level + 1)
                        lines.extend(sub_lines)

            return lines

        # 遍历所有工作表
        if isinstance(xmind_data, list):
            for sheet_idx, sheet in enumerate(xmind_data):
                if 'topic' in sheet:
                    # 添加工作表标题
                    if 'title' in sheet:
                        markdown_lines.append(f"\n# {sheet['title']}\n")
                    else:
                        markdown_lines.append(f"\n# 工作表 {sheet_idx + 1}\n")

                    # 提取主题内容
                    topic_lines = extract_topic_recursive(sheet['topic'], 0)
                    markdown_lines.extend(topic_lines)
        elif isinstance(xmind_data, dict) and 'topic' in xmind_data:
            topic_lines = extract_topic_recursive(xmind_data['topic'], 0)
            markdown_lines.extend(topic_lines)

        # 如果还有剩余图片，追加到末尾
        while image_idx < len(image_paths):
            markdown_lines.append(f"\n![placeholder]({image_paths[image_idx]})\n")
            image_idx += 1

        return "".join(markdown_lines)

    except ImportError:
        return "错误：需要安装 xmindparser 库来读取XMind文件: pip install xmindparser"
    except Exception as e:
        return f"转换XMind时出错: {e}"


def convert_to_markdown_with_placeholders(file_path: str, image_paths: List[str], temp_manager: TempFileManager) -> str:
    """
    将文档转换为带占位符的Markdown。
    """
    _, extension = os.path.splitext(file_path.lower())

    if extension == '.docx':
        return convert_docx_to_markdown_with_placeholders(file_path, image_paths, temp_manager)
    elif extension == '.pdf':
        return convert_pdf_to_markdown_with_placeholders(file_path, image_paths)
    elif extension == '.pptx':
        return convert_pptx_to_markdown_with_placeholders(file_path, image_paths, temp_manager)
    elif extension == '.xmind':
        return convert_xmind_to_markdown_with_placeholders(file_path, image_paths, temp_manager)
    else:
        # 对于其他类型，使用原始文本（暂时不支持图片占位符）
        return get_content_from_file(file_path)

# 这是分发字典，它将文件扩展名映射到正确的读取函数。
FILE_READERS = {
    '.txt': read_txt_content,
    '.docx': read_docx_content,
    '.xlsx': read_xlsx_content,
    '.pptx': read_pptx_content,
    '.pdf': read_pdf_content,
    '.xmind': read_xmind_content,
    # 在这里添加新的读取函数，例如: '.csv': read_csv_content
}


# --- 多模态LLM调用功能 ---
def encode_image_to_base64(image_path: str) -> str:
    """
    将图片文件编码为base64字符串。
    """
    try:
        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            return encoded_string
    except Exception as e:
        print(f"编码图片时出错 {image_path}: {e}")
        return ""


def analyze_images_with_qwen_vl(image_paths: List[str]) -> Dict[str, str]:
    """
    使用qwen-vl模型分析图片并返回描述结果。
    返回字典: {image_path: description}
    策略：为每张图片单独调用LLM，确保每张图片都能正确解析
    """
    try:
        # 检查API配置
        if QWEN_VL_CONFIG["api_key"] == "YOUR_API_KEY_HERE" or not QWEN_VL_CONFIG["api_key"]:
            print("警告：请先配置QWEN_VL_CONFIG中的API密钥")
            return {}

        # 初始化OpenAI客户端（使用通义千问的base_url）
        client = OpenAI(
            api_key=QWEN_VL_CONFIG["api_key"],
            base_url=QWEN_VL_CONFIG["base_url"]
        )

        image_descriptions = {}

        print(f"开始分析 {len(image_paths)} 张图片...")

        # 为每张图片单独调用LLM，确保准确性
        for idx, img_path in enumerate(image_paths, 1):
            print(f" [LLM] 正在分析图片 {idx}/{len(image_paths)}: {os.path.basename(img_path)}")

            try:
                # 编码图片
                base64_img = encode_image_to_base64(img_path)
                if not base64_img:
                    print(f" [X] 编码失败")
                    image_descriptions[img_path] = "[图片编码失败]"
                    continue

                # 构建单张图片的分析请求
                content = [
                    {
                        "type": "text",
                        "text": "请详细描述这张图片的内容，包括文字、图表、布局等所有可见信息。请用中文回答。"
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_img}"
                        }
                    }
                ]

                # 调用qwen-vl模型
                response = client.chat.completions.create(
                    model=QWEN_VL_CONFIG["model"],
                    messages=[
                        {
                            "role": "user",
                            "content": content
                        }
                    ],
                    max_tokens=1500
                )

                # 获取响应
                response_text = response.choices[0].message.content
                image_descriptions[img_path] = response_text.strip()

                # 显示描述长度作为成功标志
                desc_len = len(response_text)
                print(f" [LLM] 分析完成 (描述长度: {desc_len} 字符)")

            except Exception as e:
                error_msg = f"[图片分析失败: {str(e)}]"
                print(f" [LLM] 分析失败: {str(e)[:50]}...")
                image_descriptions[img_path] = error_msg

        print(f"图片分析完成！成功分析 {len([v for v in image_descriptions.values() if not v.startswith('[')])} / {len(image_paths)} 张图片")
        return image_descriptions

    except Exception as e:
        print(f" [LLM] 分析图片时出错: {e}")
        return {}


# --- 占位符替换功能 ---
def replace_placeholders(markdown_text: str, image_descriptions: Dict[str, str]) -> str:
    """
    将Markdown中的图片占位符替换为实际的图片描述。
    """
    try:
        # 使用正则表达式匹配 ![placeholder](image_path) 格式
        placeholder_pattern = r'!\[placeholder\]\(([^)]+)\)'

        def replace_match(match):
            image_path = match.group(1)
            # 查找对应的描述
            if image_path in image_descriptions:
                description = image_descriptions[image_path]
                # 格式化为Markdown代码块，添加长横线分隔符
                return f"\n================\n**图片描述:**\n{description}\n================\n"
            else:
                return f"\n================\n[未找到图片 {image_path} 的描述]\n================\n"

        # 执行替换
        result = re.sub(placeholder_pattern, replace_match, markdown_text)
        return result

    except Exception as e:
        print(f"替换占位符时出错: {e}")
        return markdown_text

def get_content_from_file(file_path: str) -> str:
    """
    从文件中获取内容的通用函数。
    它使用 FILE_READERS 字典来查找并调用正确的读取器。
    """
    if not os.path.exists(file_path):
        return f"错误：链接的文件 '{file_path}' 不存在"
    
    # 获取文件的扩展名
    _, extension = os.path.splitext(file_path)
    
    # 在我们的字典中查找对应的读取函数
    reader_func = FILE_READERS.get(extension.lower())
    
    if reader_func:
        # 如果找到了读取函数，就调用它
        return reader_func(file_path)
    else:
        # 否则，返回不支持的类型错误
        return f"错误：文件 '{file_path}' 的类型 ({extension}) 不受支持"

def format_as_markdown(content: str, file_extension: str) -> str:
    """
    将提取的文本内容格式化为 Markdown 代码块。
    :param content: 从文件中读取的原始文本内容。
    :param file_extension: 文件的扩展名（例如 '.txt'），用于代码块的语言标识。
    :return: 格式化后的 Markdown 字符串。
    """
    # 移除扩展名前的点，使其成为一个更干净的语言标识符
    lang_identifier = file_extension.lstrip('.')
    
    # 对于已知不支持的标识符或空标识符，使用 'text' 作为默认
    if not lang_identifier or lang_identifier in ['docx']:
        lang_identifier = 'text'
        
    return f"```{lang_identifier}\n{content}\n```"

# --- 主 Excel 处理逻辑 ---

def process_excel_in_place(excel_path: str):
    """
    自动查找链接列，在其后插入一个新列，
    用链接文档的内容填充它，并直接在原文件上保存更改。
    新版本支持图片提取和多模态LLM分析。
    """
    try:
        workbook = openpyxl.load_workbook(excel_path)
        sheet = workbook.active
        print(f"成功加载文件: '{excel_path}'")
    except FileNotFoundError:
        print(f"错误：Excel 文件 '{excel_path}' 不存在。请检查路径是否正确。")
        return
    except Exception as e:
        print(f"加载 Excel 文件 '{excel_path}' 时出错: {e}")
        return

    # 获取Excel文件所在的绝对目录
    excel_base_dir = os.path.dirname(os.path.abspath(excel_path))
    print(f"将基于此目录解析相对路径: '{excel_base_dir}'")

    all_links = [{'cell': cell, 'target': cell.hyperlink.target}
                 for row in sheet.iter_rows() for cell in row if cell.hyperlink]

    if not all_links:
        print("在此文件中未找到任何超链接。未做任何更改。")
        return

    print(f"找到了 {len(all_links)} 个超链接。")

    first_link_col_idx = all_links[0]['cell'].column
    content_col_idx = first_link_col_idx + 1

    print(f"检测到链接列为 {get_column_letter(first_link_col_idx)} 列。 "
          f"将在 {get_column_letter(content_col_idx)} 列插入新内容。")

    sheet.insert_cols(content_col_idx)

    header_cell = sheet.cell(row=1, column=content_col_idx)
    header_cell.value = "链接文档内容"
    header_cell.font = openpyxl.styles.Font(bold=True)

    # 使用临时文件管理器来管理提取的图片
    with TempFileManager() as temp_manager:
        for link_info in all_links:
            link_cell = link_info['cell']
            # 这是从Excel中读取的原始路径，可能是相对的
            relative_or_absolute_path = link_info['target']

            # 解析路径，将相对路径转换为绝对路径
            if os.path.isabs(relative_or_absolute_path):
                # 如果路径已经是绝对路径 (例如 "C:\...")，则直接使用
                full_path = relative_or_absolute_path
            else:
                # 如果是相对路径，则与Excel文件所在目录进行拼接
                full_path = os.path.join(excel_base_dir, relative_or_absolute_path)

            print(f"  - 正在处理 {link_cell.coordinate}: '{relative_or_absolute_path}' -> 解析为 '{full_path}'")

            try:
                # 步骤1: 从文档中提取图片
                print(f"    提取图片中...")
                image_paths = extract_images_from_document(full_path, temp_manager)

                if image_paths:
                    print(f"    提取到 {len(image_paths)} 张图片")
                else:
                    print(f"    未检测到图片")

                # 步骤2: 转换为带占位符的Markdown
                print(f"    转换为Markdown格式...")
                markdown_with_placeholders = convert_to_markdown_with_placeholders(
                    full_path, image_paths, temp_manager
                )

                # 步骤3: 使用LLM分析图片
                final_markdown = markdown_with_placeholders
                if image_paths:
                    print(f"    使用多模态LLM分析图片...")
                    image_descriptions = analyze_images_with_qwen_vl(image_paths)

                    if image_descriptions:
                        print(f"    替换占位符...")
                        # 步骤4: 替换占位符
                        final_markdown = replace_placeholders(
                            markdown_with_placeholders, image_descriptions
                        )
                    else:
                        print(f"    图片分析失败，使用原始内容")

                # 步骤5: 插入到Excel单元格
                content_cell = sheet.cell(row=link_cell.row, column=content_col_idx)
                content_cell.value = final_markdown

                print(f"    完成")

            except Exception as e:
                print(f"    处理出错: {e}")
                # 出错时使用原始文本
                raw_content = get_content_from_file(full_path)
                _, extension = os.path.splitext(full_path)
                md_content = format_as_markdown(raw_content, extension)
                content_cell = sheet.cell(row=link_cell.row, column=content_col_idx)
                content_cell.value = md_content

    try:
        print(f"\n正在将更改保存到原始文件: '{excel_path}'...")
        workbook.save(excel_path)
        print("处理完成！原始文件已更新。")
    except PermissionError:
        print(f"\n错误：无法保存文件。请确保 '{excel_path}' 没有被其他程序（如Excel）打开。")
    except Exception as e:
        print(f"\n保存文件 '{excel_path}' 时发生未知错误: {e}")

# --- 脚本主入口 ---
if __name__ == "__main__":
    # --- 警告 ---
    # 此脚本将直接修改您的原始文件。
    # 强烈建议在运行前对您的 Excel 文件进行备份。
    
    # --- 请在这里提供您的 Excel 文件的完整路径 ---
    excel_file_path = "C:\\Users\\Admin\\Desktop\\text\\任务管理.xlsx"
    
    process_excel_in_place(excel_file_path)